import threading
import queue
import pytesseract
import os
import re
import json
import logging
import zipfile
import fulltext
import xlrd
import hashlib
import csv
#import docx

import tensorflow as tf
#import tensorflow_hub as hub
import numpy as np
import subprocess as sp

os.environ['TF_CPP_MIN_LOG_LEVEL'] = "2"
#tf.get_logger().setLevel("ERROR")

#from official.nlp.bert import tokenization

from sklearn.feature_extraction.text import TfidfVectorizer
from io import BytesIO, StringIO, TextIOWrapper
from itertools import chain
from PIL import Image
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from gensim import corpora, models
from gensim.models import KeyedVectors, Word2Vec, ldamodel
from markdown import markdown
from xlrd.sheet import ctype_text
from hashlib import sha256

from nltk.tokenize import RegexpTokenizer
from stop_words import get_stop_words
from nltk.stem.snowball import SnowballStemmer
from nltk import ne_chunk, pos_tag, word_tokenize
from nltk.tree import Tree
from nltk.metrics.distance import edit_distance

from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfparser import PDFParser
from pdfminer.converter import HTMLConverter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfdocument import PDFDocument

from pptx import Presentation

from xml.dom.minidom import parseString
try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML

from memory_profiler import profile

from felfinder.utils import grouper, get_rand_id
from felfinder.config import W2V_MODEL_PATH

ENCODING = 'utf-8'
MAX_NUM_CHARS = 10000
MAX_GROUP_SIZE = 100
MAX_SEQ_LEN = 64

def text_proc_celery(fwrap):

    try:
        fwrap.text = process_text(fwrap.fb, fwrap.ext)
    except Exception as e:
        print("Exception in text processing thread {} on {}".format(e, fwrap.id))


class ProcessTextThread(threading.Thread):

    def __init__(self, uid, Q):
        logging.info("Initializing thread...")
        threading.Thread.__init__(self)
        self.uid = uid
        self.Q = Q
        logging.info("Thread initialized...")

    def run(self):
        logging.debug("Thread running")
        fwrap = self.Q.get()
        while (fwrap is not None):

            # try:
            #     fwrap.text = process_text(fwrap.fb, fwrap.ext)
            # except Exception as e:
            #     print("Exception in text processing thread {} on {}".format(e, fwrap.id))

            fwrap.text = process_text(fwrap.fb, fwrap.ext)


            self.Q.task_done()
            print("Finished processing id {0}".format(fwrap.id))
            fwrap = self.Q.get()

        self.Q.task_done()
        logging.info("Thread exiting")
        return


def process_text(file_bytes, ext):

    image = set(['jpg', 'jpeg', 'png', 'tiff', 'tif', 'gif', 'bmp'])
    video = set(['3gp', '3g2', 'avi', 'f4v', 'flv', 'm4v', 'asf', 'wmv', 'mpeg', 'mp4', 'qt'])
    document = set(['txt', 'rtf', 'dotx', 'dot', 'odt', 'pages', 'tex',
                    'pdf', 'ps', 'eps', 'prn', 'md', 'py', 'java', 'scala'])
    open_office = set(['odt', 'ott', 'odm', 'oth', 'ods', 'ots', 'odg',
                       'otg', 'odp', 'otp', 'odf', 'odb', 'odp'])
    doc_x = set(['docx', 'doc'])
    web = set(['html', 'xhtml', 'php', 'js', 'xml', 'war', 'ear' 'dhtml', 'mhtml'])
    spreathseet =  set(['xls', 'xlsx', 'xltx', 'xlt', 'ods', 'xlsb', 'xlsm', 'xltm'])
    presentation = set(['ppt', 'pptx', 'pot', 'potx', 'ppsx',
                        'pps', 'pptm', 'potm', 'ppsm', 'key'])

    text = None

    if ext is None:
        return ""

    if ext == "pdf":
        logging.disable(logging.CRITICAL)
        text = pdf_text_encode(file_bytes, ENCODING)
        logging.disable(logging.NOTSET)
    elif ext == "csv":
        text = csv_text_encode(file_bytes, ENCODING)
    elif ext == "tsv":
        text = csv_text_encode(file_bytes, ENCODING, "\t")
    elif ext == "doc":
        text = doc_text_encode(file_bytes, ENCODING)
    elif ext == "rtf":
        text = None #general_text_extract(file_bytes, ENCODING, "rtf")
    elif ext == "md":
        text = md_text_encode(file_bytes, ENCODING)
    elif ext == "html":
        text = general_text_extract(file_bytes, ENCODING, "html")
    elif ext == "json":
        text = None #general_text_extract(file_bytes, ENCODING, "json")
    elif ext in doc_x:
        text = docx_text_encode(file_bytes, ENCODING)
    elif ext in open_office:
        text = open_office_text_encode(file_bytes, ENCODING)
    elif ext in document:
        text = pure_text_encode(file_bytes, ENCODING)
    elif ext in image:
        text = img_text_encode(file_bytes, ENCODING)
    elif ext in spreathseet:
        text = spreadsheet_text_encode(file_bytes, ENCODING)
    elif ext in presentation:
        text = pptx_text_encode(file_bytes, ENCODING)

    if text is None:
        print("Unprocessing extension: {0}".format(ext))

    if text is not None:
        #print("{0}: {1}".format(ext, text[:100]))
        return text
    else:
        return ""

def get_bert_tokenizer(bert_layer):

    vocab_file = bert_layer.resolved_object.vocab_file.asset_path.numpy()
    do_lower_case = bert_layer.resolved_object.do_lower_case.numpy()
    tokenizer = tokenization.FullTokenizer(vocab_file, do_lower_case)

    return tokenizer

# def get_bert_layer(HUB_URL, max_seq_length):


#     in_shape = tf.convert_to_tensor((max_seq_length,))
#     input_word_ids = tf.keras.layers.Input(shape=in_shape, dtype=tf.int32,
#                                            name="input_word_ids")
#     input_mask = tf.keras.layers.Input(shape=in_shape, dtype=tf.int32,
#                                        name="input_mask")
#     segment_ids = tf.keras.layers.Input(shape=in_shape, dtype=tf.int32,
#                                         name="segment_ids")
#     bert_layer = hub.KerasLayer(HUB_URL,
#                                 trainable=False)
#     return bert_layer

def bert_input(text, tokenizer, max_seq_length):

  tokens_a = tokenizer.tokenize(text)
  # Account for [CLS] and [SEP] with "- 2"
  if len(tokens_a) > max_seq_length - 2:
      tokens_a = tokens_a[0:(max_seq_length - 2)]

  # The convention in BERT is:
  # (a) For sequence pairs:
  #  tokens:   [CLS] is this jack ##son ##ville ? [SEP] no it is not . [SEP]
  #  type_ids: 0     0  0    0    0     0       0 0     1  1  1  1   1 1
  # (b) For single sequences:
  #  tokens:   [CLS] the dog is hairy . [SEP]
  #  type_ids: 0     0   0   0  0     0 0
  #
  # Where "type_ids" are used to indicate whether this is the first
  # sequence or the second sequence. The embedding vectors for `type=0` and
  # `type=1` were learned during pre-training and are added to the wordpiece
  # embedding vector (and position vector). This is not *strictly* necessary
  # since the [SEP] token unambiguously separates the sequences, but it makes
  # it easier for the model to learn the concept of sequences.
  #
  # For classification tasks, the first vector (corresponding to [CLS]) is
  # used as the "sentence vector". Note that this only makes sense because
  # the entire model is fine-tuned.
  tokens = []
  segment_ids = []
  tokens.append("[CLS]")
  segment_ids.append(0)
  for token in tokens_a:
    tokens.append(token)
    segment_ids.append(0)
  tokens.append("[SEP]")
  segment_ids.append(0)

  input_ids = tokenizer.convert_tokens_to_ids(tokens)

  # The mask has 1 for real tokens and 0 for padding tokens. Only real
  # tokens are attended to.
  input_mask = [1] * len(input_ids)

  # Zero-pad up to the sequence length.
  while len(input_ids) < max_seq_length:
    input_ids.append(0)
    input_mask.append(0)
    segment_ids.append(0)

  assert len(input_ids) == max_seq_length
  assert len(input_mask) == max_seq_length
  assert len(segment_ids) == max_seq_length

  return input_ids, input_mask, segment_ids

def finalize_text_rep(all_wrappers):

    HUB_URL = "https://tfhub.dev/tensorflow/bert_en_cased_L-12_H-768_A-12/1"
    max_seq_length = MAX_SEQ_LEN  # Your choice here.

    bert_layer = get_bert_layer(HUB_URL, tf.constant(max_seq_length))
    bert_token = get_bert_tokenizer(bert_layer)

    for fwrap in all_wrappers:
        if fwrap.usable_text:

            mod_txt = fwrap.text[:MAX_NUM_CHARS]
            mod_txt = mod_txt.replace('\n', '.')
            sentences = mod_txt.split('.')

            full_out = []

            input_ids, input_mask, segment_ids = [], [], []
            for s_group in grouper(sentences, MAX_GROUP_SIZE):
                for s in s_group:
                    if not s: #Due to grouper, s could be None
                        continue
                    if len(s) >= max_seq_length:
                        continue

                    iid, mask, seg = bert_input(s, bert_token, max_seq_length)
                    input_ids.append(iid)
                    input_mask.append(mask)
                    segment_ids.append(seg)

                if input_ids:
                    temp_out, _ = bert_layer([tf.convert_to_tensor(input_ids),
                                              tf.convert_to_tensor(input_mask),
                                              tf.convert_to_tensor(segment_ids)])
                    full_out.append(temp_out)

            if full_out:
                full_out = np.vstack(full_out)
                full_out = np.mean(full_out, axis=0)
                fwrap.text_rep = full_out

def tfidf_analyze(doc_list):

    tfidf_vect = TfidfVectorizer(max_df = .8, min_df = .2, stop_words = "english", \
                                 use_idf = True, tokenizer=tok_and_stem)
    tfidf_matrix = tfidf_vect.fit_transform(doc_list)

    return tfidf_matrix, tfidf_vect.vocabulary_


def edit_dist(s1, s2):
    if s1 is None or s2 is None:
        return 1000
    return edit_distance(s1, s2)

def doc_vec(model, doc):
    """
    Create a makeshift document vector by taking the means of all word vectors
    """

    tok_and_stemmed_doc = tok_and_stem(doc)
    clean_doc = [word for word in tok_and_stemmed_doc if word in model.vocab]
    if clean_doc:
        return np.mean(model[clean_doc], axis = 0)
    else:
        example_vector = model.get_vector("garbage")
        return np.zeros(example_vector.shape)

def load_model():
    """
    Load the word2vec pretrained model from the data file
    """

    WORD_VECTORS_TO_READ = 500000

    model_fname = W2V_MODEL_PATH
    model = KeyedVectors.load_word2vec_format(model_fname, binary = True, \
                                              limit = WORD_VECTORS_TO_READ)

    model.init_sims(replace = True)

    return model

def word_to_vec(doc_list):
    """
    Returns the similarity between all items
    """

    model = load_model()
    corpus = [doc_vec(model, doc) for doc in doc_list]

    return corpus


def tok_and_stem(doc):

    tokenizer = RegexpTokenizer(r'\w+')
    en_stop = get_stop_words('en')
    snowball_stemmer = SnowballStemmer('english')

    raw = doc.lower()
    tokens = tokenizer.tokenize(raw)

    stopped_tokens = [it for it in tokens if not it in en_stop]
    stemmed_tokens = [snowball_stemmer.stem(it) for it in stopped_tokens]
    return stemmed_tokens


def pure_text_encode(f_bytes, encoding):
    return f_bytes.decode(encoding)

def csv_text_encode(f_bytes, encoding, delim = ","):
    return f_bytes.decode(encoding).replace(delim, " ")

def get_metadata(fb, ext):

    if not ext:
        return ""

    if ext.lower() == 'pdf':
        logging.disable(logging.CRITICAL)
        meta = get_pdf_metadata(fb)
        logging.disable(logging.NOTSET)
        return meta
    elif ext.lower() == 'docx':
        return ""
        #return get_docx_metadata(fb)
    else:
        return ""

def get_pdf_metadata(fb):

    doc = PDFDocument(PDFParser(BytesIO(fb)))
    info = doc.info[0]

    for k in info:
        if type(info[k]) == bytes:
            try:

                if k.lower() == 'author' or k.lower() == 'creator':
                    sh = sha256()
                    sh.update(info[k])
                    info[k] = sh.hexdigest()
                else:
                    info[k] = info[k].decode('utf-8')

            except UnicodeDecodeError:
                info[k] = ""

    try:
        js = json.dumps(info)
    except:
        js = ""

    return js

def orig_pdf_text_encode(f_bytes, encoding):

    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    laparams = LAParams()

    device = TextConverter(rsrcmgr, retstr, codec=encoding, laparams=laparams)
    fp = BytesIO(f_bytes)
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    maxpages = 0
    pagenos=set()

    for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages,
                                  password="",caching=True, check_extractable=True):
        interpreter.process_page(page)

    text = retstr.getvalue()
    device.close()
    retstr.close()

    return text


def pdf_text_encode(f_bytes, encoding):

    base_id = get_rand_id()
    fake_fname = base_id + ".pdf"
    with open(fake_fname, "wb") as of:
        of.write(f_bytes)

    stdout = sp.run(["pdftotext", fake_fname])

    with open(base_id + '.txt', 'r') as of:
        text = of.read()
    os.system("shred -u " + fake_fname)
    os.system("shred -u " + base_id + '.txt')

    return text


def img_text_encode(f_bytes, encoding):
    im_to_txt = Image.open(BytesIO(f_bytes))
    return pytesseract.image_to_string(img)

def xml_text_encode(f_bytes, encoding):

    doc = parseString(f_bytes)
    paragraphs = doc.getElementsByTagName('text:p')

    text = [str(ch.data) for ch in filter(\
                                     lambda x: x.nodeType == x.TEXT_NODE, \
                                     chain(*[p.childNodes for p in paragraphs]))]
    return " ".join(text)

def open_office_text_encode(f_bytes, encoding):

    """
    Open office files are essentially zipped archives. The key file
    is the content.xml file within the archive, which can then
    be parsed to extract the text.
    """

    open_office_file = zipfile.ZipFile(BytesIO(f_bytes))
    return xml_text_encode(open_office_file.read('content.xml'), encoding)


#This could be improved...
def spreadsheet_text_encode(f_bytes, encoding):

    #UTF-8 is assumed for encoding, which isn't great. May want to modify later.

    wb = xlrd.open_workbook(file_contents = f_bytes)
    text = []
    for sheet in wb.sheets():
        for row in sheet.get_rows():
            filtered_row = filter(lambda x: ctype_text.get(x.ctype, 'not_text') == 'text', row)
            filtered_row = [s.value for s in filtered_row]
            text += [" ".join(filtered_row)]
    return " ".join(text)

# def get_docx_metadata(fb):

#     """
#     Returns a JSON string of the available information
#     """

#     metadata = {}
#     doc = docx.Document(BytesIO(fb))

#     prop = doc.core_properties

#     # Hash author name if it exists
#     if prop.author != '':
#         sh = sha256()
#         sh.update(prop.author.encode('utf-8'))
#         metadata["author"] = sh.hexdigest()

#     metadata["category"] = prop.category
#     metadata["comments"] = prop.comments
#     metadata["content_status"] = prop.content_status
#     metadata["created"] = prop.created
#     metadata["identifier"] = prop.identifier
#     metadata["keywords"] = prop.keywords
#     metadata["language"] = prop.language
#     metadata["modified"] = prop.modified
#     metadata["subject"] = prop.subject
#     metadata["title"] = prop.title
#     metadata["version"] = prop.version

#     if metadata['created'] is not None and metadata["created"] != '' and \
#        not type(metadata["created"]) == str:
#         metadata["created"] = metadata["created"].strftime("%Y-%m-%dT%H:%M:%SZ")

#     if metadata['modified'] is not None and metadata["modified"] != '' and \
#        not type(metadata["modified"]) == str:
#         metadata["modified"] = metadata["modified"].strftime("%Y-%m-%dT%H:%M:%SZ")

#     try:
#         js = json.dumps(metadata)
#     except:
#         js = ""

#     return js

#Reference: https://etienned.github.io/posts/extract-text-from-word-docx-simply/
def docx_text_encode(f_bytes, encoding):
    """
    Take the path of a docx file as argument, return the text in unicode.
    """

    WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
    PARA = WORD_NAMESPACE + 'p'
    TEXT = WORD_NAMESPACE + 't'

    document = zipfile.ZipFile(BytesIO(f_bytes))
    xml_content = document.read('word/document.xml')

    document.close()
    tree = XML(xml_content)

    paragraphs = []
    for paragraph in tree.getiterator(PARA):
        texts = [node.text
                 for node in paragraph.getiterator(TEXT)
                 if node.text]
        if texts:
            paragraphs.append(''.join(texts))

    return '\n\n'.join(paragraphs)

def general_text_extract(f_bytes, encoding, ext):

    f_type = "_." + ext
    return fulltext.get(BytesIO(f_bytes), name = f_type)

def md_text_encode(f_bytes, encoding):
    html = markdown(f_bytes.decode('utf-8'))
    return fulltext.get(StringIO(html), name = "_.html")

def doc_text_encode(f_bytes, encoding):

    fake_fname = get_rand_id() + ".doc"
    with open(fake_fname, "wb") as of:
        of.write(f_bytes)

    text = os.popen("antiword " + fake_fname).read()
    os.system("shred -u " + fake_fname)

    return text

def pptx_text_encode(f_bytes, encoding):

    prs = Presentation(BytesIO(f_bytes))

    text = ""

    for sld in prs.slides:
        for shape in sld.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                text += str(p.text) + "\n"

    return text
